from __future__ import annotations

import io
from pathlib import Path

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".pdf", ".docx", ".xlsx", ".csv", ".pptx", "text/plain", "text/csv",
}


def parse_bytes(name: str, content: bytes, mime_type: str) -> str:
    ext = Path(name).suffix.lower()
    if mime_type == "text/plain" or ext in (".txt", ".md"):
        return content.decode("utf-8", errors="replace")
    if mime_type == "text/csv" or ext == ".csv":
        return content.decode("utf-8", errors="replace")
    if ext == ".pdf":
        return _parse_pdf(content)
    if ext == ".docx":
        return _parse_docx(content)
    if ext == ".xlsx":
        return _parse_xlsx(content)
    if ext == ".pptx":
        return _parse_pptx(content)
    return content.decode("utf-8", errors="replace")


def _parse_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    return "\n".join(p.extract_text() or "" for p in reader.pages)


def _parse_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def _parse_xlsx(content: bytes) -> str:
    import pandas as pd
    dfs = pd.read_excel(io.BytesIO(content), sheet_name=None)
    return "\n\n".join(
        f"[{sheet}]\n{df.to_string(index=False)}" for sheet, df in dfs.items()
    )


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
    return "\n".join(lines)
